import os
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
from PIL import Image
import pytesseract
from deep_translator import GoogleTranslator
import PyPDF2
from docx import Document
from pptx import Presentation
import reportlab.lib.styles
from reportlab.pdfgen import canvas

# Configure Tesseract path (change according to your installation)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

SUPPORTED_LANGUAGES = {
    'auto': 'Auto Detect',
    'af': 'Afrikaans',
    'sq': 'Albanian',
    'am': 'Amharic',
    'ar': 'Arabic',
    'hy': 'Armenian',
    'az': 'Azerbaijani',
    'eu': 'Basque',
    'be': 'Belarusian',
    'bn': 'Bengali',
    'bs': 'Bosnian',
    'bg': 'Bulgarian',
    'my': 'Burmese',
    'ca': 'Catalan',
    'ceb': 'Cebuano',
    'zh-CN': 'Chinese (Simplified)',
    'zh-TW': 'Chinese (Traditional)',
    'co': 'Corsican',
    'hr': 'Croatian',
    'cs': 'Czech',
    'da': 'Danish',
    'nl': 'Dutch',
    'en': 'English',
    'eo': 'Esperanto',
    'et': 'Estonian',
    'fi': 'Finnish',
    'fr': 'French',
    'fy': 'Frisian',
    'gl': 'Galician',
    'ka': 'Georgian',
    'de': 'German',
    'el': 'Greek',
    'gu': 'Gujarati',
    'ht': 'Haitian Creole',
    'ha': 'Hausa',
    'haw': 'Hawaiian',
    'he': 'Hebrew',
    'hi': 'Hindi',
    'hmn': 'Hmong',
    'hu': 'Hungarian',
    'is': 'Icelandic',
    'ig': 'Igbo',
    'id': 'Indonesian',
    'ga': 'Irish',
    'it': 'Italian',
    'ja': 'Japanese',
    'jw': 'Javanese',
    'kn': 'Kannada',
    'kk': 'Kazakh',
    'km': 'Khmer',
    'rw': 'Kinyarwanda',
    'ko': 'Korean',
    'ku': 'Kurdish',
    'ky': 'Kyrgyz',
    'lo': 'Lao',
    'la': 'Latin',
    'lv': 'Latvian',
    'lt': 'Lithuanian',
    'lb': 'Luxembourgish',
    'mk': 'Macedonian',
    'mg': 'Malagasy',
    'ms': 'Malay',
    'ml': 'Malayalam',
    'mt': 'Maltese',
    'mi': 'Maori',
    'mr': 'Marathi',
    'mn': 'Mongolian',
    'ne': 'Nepali',
    'no': 'Norwegian',
    'ny': 'Nyanja',
    'or': 'Odia',
    'ps': 'Pashto',
    'fa': 'Persian',
    'pl': 'Polish',
    'pt': 'Portuguese',
    'pa': 'Punjabi',
    'ro': 'Romanian',
    'ru': 'Russian',
    'sm': 'Samoan',
    'gd': 'Scots Gaelic',
    'sr': 'Serbian',
    'st': 'Sesotho',
    'sn': 'Shona',
    'sd': 'Sindhi',
    'si': 'Sinhala',
    'sk': 'Slovak',
    'sl': 'Slovenian',
    'so': 'Somali',
    'es': 'Spanish',
    'su': 'Sundanese',
    'sw': 'Swahili',
    'sv': 'Swedish',
    'tl': 'Tagalog',
    'tg': 'Tajik',
    'ta': 'Tamil',
    'tt': 'Tatar',
    'te': 'Telugu',
    'th': 'Thai',
    'tr': 'Turkish',
    'tk': 'Turkmen',
    'uk': 'Ukrainian',
    'ur': 'Urdu',
    'ug': 'Uyghur',
    'uz': 'Uzbek',
    'vi': 'Vietnamese',
    'cy': 'Welsh',
    'xh': 'Xhosa',
    'yi': 'Yiddish',
    'yo': 'Yoruba',
    'zu': 'Zulu',
    # Additional languages
    'ak': 'Akan',
    'as': 'Assamese',
    'ay': 'Aymara',
    'bm': 'Bambara',
    'bi': 'Bislama',
    'dv': 'Divehi',
    'dz': 'Dzongkha',
    'ee': 'Ewe',
    'fo': 'Faroese',
    'fj': 'Fijian',
    'gn': 'Guarani',
    'ia': 'Interlingua',
    'ie': 'Interlingue',
    'iu': 'Inuktitut',
    'ik': 'Inupiaq',
    'kl': 'Kalaallisut',
    'ki': 'Kikuyu',
    'rn': 'Kirundi',
    'kg': 'Kongo',
    'kj': 'Kuanyama',
    'lg': 'Ganda',
    'ln': 'Lingala',
    'lu': 'Luba-Katanga',
    'mh': 'Marshallese',
    'na': 'Nauru',
    'ng': 'Ndonga',
    'oc': 'Occitan',
    'om': 'Oromo',
    'pi': 'Pali',
    'qu': 'Quechua',
    'sa': 'Sanskrit',
    'sc': 'Sardinian',
    'sg': 'Sango',
    'ss': 'Swati',
    'ti': 'Tigrinya',
    'ts': 'Tsonga',
    'tn': 'Tswana',
    'to': 'Tonga',
    'ty': 'Tahitian',
    've': 'Venda',
    'vo': 'Volapük',
    'wa': 'Walloon',
    'wo': 'Wolof',
    'za': 'Zhuang'
}

# Expanded Tesseract language code exceptions
TESSERACT_EXCEPTIONS = {
    'zh-CN': 'chi_sim',
    'zh-TW': 'chi_tra',
    'hi': 'hin',
    'bn': 'ben',
    'ja': 'jpn',
    'ko': 'kor',
    'te': 'tel',
    'kn': 'kan',
    'ml': 'mal',
    'ta': 'tam',
    'gu': 'guj',
    'mr': 'mar',
    'pa': 'pan',
    'ne': 'nep',
    'or': 'ori',
    'as': 'asm',
    'sa': 'san',
    'si': 'sin',
    'ur': 'urd',
    'ar': 'ara',
    'fa': 'fas',
    'he': 'heb',
    'am': 'amh',
    'ti': 'tir',
    'km': 'khm',
    'th': 'tha',
    'my': 'mya',
    'bo': 'bod',
    'dz': 'dzo',
    'ka': 'kat',
    'hy': 'hye',
    'ru': 'rus',
    'uk': 'ukr',
    'be': 'bel',
    'bg': 'bul',
    'mk': 'mkd',
    'sr': 'srp',
    'hr': 'hrv',
    'bs': 'bos',
    'sl': 'slv',
    'cs': 'ces',
    'sk': 'slk',
    'pl': 'pol',
    'hu': 'hun',
    'ro': 'ron',
    'el': 'ell',
    'tr': 'tur',
    'az': 'aze',
    'kk': 'kaz',
    'uz': 'uzb',
    'ky': 'kir',
    'tg': 'tgk',
    'mn': 'mon',
    'ps': 'pus',
    'ku': 'kur',
    'sd': 'snd',
    'dv': 'div',
    'yi': 'yid',
    'ceb': 'ceb',
    'fil': 'fil',
    'haw': 'haw',
    'hmn': 'hmn',
    'lo': 'lao',
    'gl': 'glg',
    'co': 'cos',
    'ht': 'hat',
    'sm': 'smo',
    'st': 'sot',
    'tn': 'tsn',
    'ts': 'tso',
    'ss': 'ssw',
    've': 'ven',
    'mi': 'mri',
    'gd': 'gla',
    'ga': 'gle',
    'cy': 'cym',
    'fy': 'fry',
    'xh': 'xho',
    'zu': 'zul',
    'sn': 'sna',
    'yo': 'yor',
    'ig': 'ibo',
    'ha': 'hau',
    'so': 'som',
    'sw': 'swa',
    'mg': 'mlg',
    'eo': 'epo',
    'la': 'lat',
    'tl': 'tgl',
    'ny': 'nya',
    'lg': 'lug',
    'rw': 'kin',
    'ak': 'aka',
    'bm': 'bam',
    'ee': 'ewe',
    'ff': 'ful',
    'kl': 'kal',
    'ki': 'kik',
    'ln': 'lin',
    'lg': 'lug',
    'rn': 'run',
    'sg': 'sag',
    'wo': 'wol',
    'xh': 'xho',
    'yo': 'yor',
    'zu': 'zul'
}

def get_tesseract_lang(google_code):
    """Convert Google Translate language code to Tesseract OCR code"""
    if google_code in TESSERACT_EXCEPTIONS:
        return TESSERACT_EXCEPTIONS[google_code]
    return google_code[:3]  # Default to first 3 letters of ISO code

class DocumentTranslatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Multi-Language Document Translator (100+ Languages)")
        self.setup_ui()
        self.files_to_process = []
        self.running = False

    def setup_ui(self):
        main_frame = ttk.Frame(self.root, padding=10)
        main_frame.grid(row=0, column=0, sticky='nsew')

        # File selection
        ttk.Label(main_frame, text="Select Files:").grid(row=0, column=0, sticky='w')
        self.file_listbox = tk.Listbox(main_frame, width=60, height=5)
        self.file_listbox.grid(row=1, column=0, columnspan=2, pady=5)
        ttk.Button(main_frame, text="Add Files", command=self.add_files).grid(row=2, column=0, pady=5)
        ttk.Button(main_frame, text="Clear List", command=self.clear_files).grid(row=2, column=1, pady=5)

        # Language selection
        ttk.Label(main_frame, text="Source Language:").grid(row=3, column=0, sticky='w', pady=5)
        self.src_lang = ttk.Combobox(main_frame, values=list(SUPPORTED_LANGUAGES.values()), state='readonly')
        self.src_lang.set('Auto Detect')
        self.src_lang.grid(row=4, column=0, sticky='ew')
        
        ttk.Label(main_frame, text="Target Language:").grid(row=3, column=1, sticky='w', pady=5)
        self.tgt_lang = ttk.Combobox(main_frame, values=list(SUPPORTED_LANGUAGES.values()), state='readonly')
        self.tgt_lang.set('English')
        self.tgt_lang.grid(row=4, column=1, sticky='ew')

        # Progress bar
        self.progress = ttk.Progressbar(main_frame, orient='horizontal', mode='determinate')
        self.progress.grid(row=5, column=0, columnspan=2, pady=10, sticky='ew')

        # Translate button
        ttk.Button(main_frame, text="Translate Documents", command=self.start_translation).grid(row=6, column=0, columnspan=2, pady=10)

        # Translated content display
        ttk.Label(main_frame, text="Translated Content:").grid(row=7, column=0, sticky='w', pady=5)
        self.translated_text_area = scrolledtext.ScrolledText(main_frame, wrap=tk.WORD, width=80, height=20)
        self.translated_text_area.grid(row=8, column=0, columnspan=2, pady=10)

        # Configure grid weights
        main_frame.columnconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)

    def add_files(self):
        filetypes = [
            ('Documents', '*.pdf *.docx *.pptx *.txt'),
            ('Images', '*.png *.jpg *.jpeg *.bmp'),
            ('All files', '*.*')
        ]
        files = filedialog.askopenfilenames(filetypes=filetypes)
        for f in files:
            self.file_listbox.insert('end', f)
            self.files_to_process.append(f)

    def clear_files(self):
        self.file_listbox.delete(0, 'end')
        self.files_to_process = []

    def start_translation(self):
        if not self.files_to_process:
            messagebox.showwarning("No Files", "Please select files to translate")
            return

        src_lang_code = self.get_lang_code(self.src_lang.get())
        tgt_lang_code = self.get_lang_code(self.tgt_lang.get())

        if not tgt_lang_code:
            messagebox.showwarning("Invalid Language", "Please select a valid target language")
            return

        self.running = True
        threading.Thread(target=self.process_files, args=(src_lang_code, tgt_lang_code)).start()

    def get_lang_code(self, lang_name):
        for code, name in SUPPORTED_LANGUAGES.items():
            if name == lang_name:
                return code
        return None

    def process_files(self, src_lang, tgt_lang):
        total_files = len(self.files_to_process)
        for i, file_path in enumerate(self.files_to_process):
            if not self.running:
                break
            try:
                self.update_progress((i/total_files)*100)
                output_path, translated_text = self.translate_file(file_path, src_lang, tgt_lang)
                self.display_translated_content(translated_text)
                messagebox.showinfo("Success", f"File translated successfully!\nSaved to: {output_path}")
            except pytesseract.TesseractError as e:
                messagebox.showerror("OCR Error", f"Please install Tesseract language pack for {SUPPORTED_LANGUAGES.get(src_lang, 'selected language')}")
            except Exception as e:
                messagebox.showerror("Error", f"Error processing {os.path.basename(file_path)}:\n{str(e)}")
        self.update_progress(0)
        self.running = False

    def translate_file(self, file_path, src_lang, tgt_lang):
        ext = os.path.splitext(file_path)[1].lower()
        handler = {
            '.pdf': self.handle_pdf,
            '.docx': self.handle_docx,
            '.pptx': self.handle_pptx,
            '.txt': self.handle_text,
            '.png': self.handle_image,
            '.jpg': self.handle_image,
            '.jpeg': self.handle_image,
            '.bmp': self.handle_image
        }.get(ext, self.handle_unknown)

        return handler(file_path, src_lang, tgt_lang)

    def handle_pdf(self, file_path, src_lang, tgt_lang):
        translated_text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                translated = self.translate_text(text, src_lang, tgt_lang)
                translated_text.append(translated)

        output_path = file_path.replace('.pdf', '_translated.pdf')
        self.create_pdf(output_path, translated_text)
        return output_path, "\n".join(translated_text)

    def handle_docx(self, file_path, src_lang, tgt_lang):
        doc = Document(file_path)
        translated_doc = Document()
        translated_text = []
        
        for para in doc.paragraphs:
            translated = self.translate_text(para.text, src_lang, tgt_lang)
            translated_doc.add_paragraph(translated)
            translated_text.append(translated)
        
        output_path = file_path.replace('.docx', '_translated.docx')
        translated_doc.save(output_path)
        return output_path, "\n".join(translated_text)

    def handle_pptx(self, file_path, src_lang, tgt_lang):
        prs = Presentation(file_path)
        translated_prs = Presentation()
        translated_text = []

        for slide in prs.slides:
            translated_slide = translated_prs.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    translated = self.translate_text(shape.text, src_lang, tgt_lang)
                    new_shape = translated_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height)
                    new_shape.text = translated
                    translated_text.append(translated)

        output_path = file_path.replace('.pptx', '_translated.pptx')
        translated_prs.save(output_path)
        return output_path, "\n".join(translated_text)

    def handle_text(self, file_path, src_lang, tgt_lang):
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        
        translated_text = self.translate_text(text, src_lang, tgt_lang)
        output_path = file_path.replace('.txt', '_translated.txt')
        with open(output_path, 'w', encoding='utf-8') as file:
            file.write(translated_text)
        return output_path, translated_text

    def handle_image(self, file_path, src_lang, tgt_lang):
        img = Image.open(file_path)
        
        # Determine Tesseract language
        tess_lang = 'eng' if src_lang == 'auto' else get_tesseract_lang(src_lang)
        
        try:
            text = pytesseract.image_to_string(img, lang=tess_lang)
        except pytesseract.TesseractError:
            text = pytesseract.image_to_string(img)  # Fallback to English
        
        translated_text = self.translate_text(text, src_lang, tgt_lang)
        
        output_path = file_path.replace(os.path.splitext(file_path)[1], '_translated.txt')
        with open(output_path, 'w', encoding='utf-8') as file:
            file.write(translated_text)
        return output_path, translated_text

    def handle_unknown(self, file_path, src_lang, tgt_lang):
        raise ValueError("Unsupported file format")

    def translate_text(self, text, src_lang, tgt_lang):
        try:
            if not text.strip():
                return ""
            # Use 'auto' as source if selected
            translated = GoogleTranslator(source=src_lang if src_lang != 'auto' else 'auto', 
                                        target=tgt_lang).translate(text)
            return translated
        except Exception as e:
            raise RuntimeError(f"Translation failed: {str(e)}")

    def create_pdf(self, output_path, pages):
        c = canvas.Canvas(output_path)
        style = reportlab.lib.styles.getSampleStyleSheet()
        width, height = (612, 792)  # Letter size
        
        for page in pages:
            text = c.beginText(40, height - 40)
            text.setFont("Helvetica", 12)
            text.textLines(page)
            c.drawText(text)
            c.showPage()
        c.save()

    def display_translated_content(self, content):
        self.translated_text_area.delete(1.0, tk.END)  # Clear previous content
        self.translated_text_area.insert(tk.END, content)  # Insert new content

    def update_progress(self, value):
        self.root.after(0, lambda: self.progress.configure(value=value))

    def on_closing(self):
        self.running = False
        self.root.destroy()

if __name__ == "__main__":
    root = tk.Tk()
    app = DocumentTranslatorApp(root)
    root.protocol("WM_DELETE_WINDOW", app.on_closing)
    root.mainloop()